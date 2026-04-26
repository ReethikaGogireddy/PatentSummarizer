"""
Generates Group22-Project9-SP26-Group-DEMO-Presentation.pptx
Run: py -3.11 make_demo_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image
import os

# Colors
NAVY   = RGBColor(0x0f, 0x34, 0x60)
BLUE   = RGBColor(0x16, 0x21, 0x3e)
GOLD   = RGBColor(0xe9, 0xc4, 0x6a)
WHITE  = RGBColor(0xff, 0xff, 0xff)
LIGHT  = RGBColor(0xf0, 0xf4, 0xff)
GREEN  = RGBColor(0x21, 0x96, 0x53)
ORANGE = RGBColor(0xe7, 0x6f, 0x51)

OUT_DIR = os.path.join(os.path.dirname(__file__), "EVALUATIONS")

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]  # completely blank


# ── helpers ──────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(blank)

def rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return box

def header_bar(slide, title, subtitle=""):
    rect(slide, 0, 0, 13.33, 1.4, fill=NAVY)
    txt(slide, title, 0.35, 0.12, 12, 0.65,
        size=28, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.35, 0.78, 12, 0.5,
            size=14, color=GOLD)

def footer(slide, text="Group 22  |  Project 9  |  CSE 573 Spring 2026  |  Arizona State University"):
    rect(slide, 0, 7.15, 13.33, 0.35, fill=BLUE)
    txt(slide, text, 0.3, 7.17, 13, 0.28,
        size=11, color=RGBColor(0xa8, 0xd8, 0xea))

def fit_picture(slide, path, l, t, max_w, max_h):
    """Add a picture to slide centered in the (l,t,max_w,max_h) box,
    preserving the source image's aspect ratio. Returns the picture shape."""
    if not os.path.exists(path):
        txt(slide, f"[Missing: {os.path.basename(path)}]",
            l, t + max_h/2 - 0.2, max_w, 0.4,
            size=14, color=NAVY, align=PP_ALIGN.CENTER, italic=True)
        return None
    img_w, img_h = Image.open(path).size
    src_ratio = img_w / img_h
    box_ratio = max_w / max_h
    if src_ratio >= box_ratio:
        # image is wider than the box -> fit width, scale height
        w = max_w
        h = max_w / src_ratio
    else:
        h = max_h
        w = max_h * src_ratio
    cx = l + (max_w - w) / 2
    cy = t + (max_h - h) / 2
    return slide.shapes.add_picture(path, Inches(cx), Inches(cy),
                                     Inches(w), Inches(h))

def bullet_block(slide, items, l, t, w, h,
                 size=16, color=RGBColor(0x22, 0x22, 0x22),
                 bullet="•"):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size  = Pt(size)
        run.font.color.rgb = color
    return box


# ── Slide 1: Cover ───────────────────────────────────────────────────────────

sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
rect(sl, 0, 0, 13.33, 0.08, fill=GOLD)
rect(sl, 0, 7.42, 13.33, 0.08, fill=GOLD)

txt(sl, "Patent Corpus Semantic Analysis",
    1, 1.5, 11.33, 1.2, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(sl, "DEMO Presentation",
    1, 2.8, 11.33, 0.7, size=26, bold=False, color=GOLD, align=PP_ALIGN.CENTER)

rect(sl, 3.5, 3.7, 6.33, 0.04, fill=GOLD)

txt(sl, "Group 22  |  Project 9",
    1, 3.9, 11.33, 0.5, size=18, color=WHITE, align=PP_ALIGN.CENTER)
txt(sl, "CSE 573 Spring 2026  |  Arizona State University",
    1, 4.45, 11.33, 0.5, size=16, color=RGBColor(0xa8, 0xd8, 0xea), align=PP_ALIGN.CENTER)
txt(sl, "Akshith Reddy Vempati  |  Reethika Gogireddy",
    1, 5.0, 11.33, 0.5, size=15, color=RGBColor(0xa8, 0xd8, 0xea), align=PP_ALIGN.CENTER)

txt(sl, "April 27, 2026",
    1, 6.3, 11.33, 0.5, size=14, color=GOLD, align=PP_ALIGN.CENTER)


# ── Slide 2: Problem Statement ───────────────────────────────────────────────

sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xf4, 0xf6, 0xf9))
header_bar(sl, "Problem Statement", "Why automated patent analysis?")
footer(sl)

items = [
    "4+ million active US patents across hundreds of CPC technology classes",
    "Manual review is impossible at scale — researchers need fast discovery tools",
    "Patents use highly technical, domain-specific language with minimal overlap",
    "Existing tools lack semantic understanding — keyword search misses synonyms and paraphrases",
    "Goal: cluster large patent corpora by topic, extract cluster summaries, and rank by relevance",
]
bullet_block(sl, items, 0.6, 1.7, 12.1, 4.5, size=17,
             color=RGBColor(0x22, 0x22, 0x22))

rect(sl, 0.6, 6.3, 12.1, 0.6, fill=NAVY)
txt(sl, "Our system automates patent discovery using TF-IDF, LDA, and Sentence-BERT embeddings "
        "combined with K-Means and Hierarchical clustering to reveal hidden topical structure.",
    0.75, 6.35, 11.8, 0.55, size=13, color=WHITE)


# ── Slide 3: System Architecture ─────────────────────────────────────────────

sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xf4, 0xf6, 0xf9))
header_bar(sl, "System Architecture", "Two-stage pipeline: representation, then analysis")
footer(sl)

# Two rows of stages with grouped visual hierarchy.
# Row 1 (Input + Representation): 3 boxes
# Row 2 (Analysis + Output):       4 boxes

LIGHT_NAVY = RGBColor(0x16, 0x40, 0x70)
DEEP_NAVY  = RGBColor(0x0a, 0x22, 0x40)

row1 = [
    ("1. Ingestion",       "Real USPTO patents\nHuggingFace big_patent\n5,000 docs across\n5 CPC domains"),
    ("2. Preprocessing",   "Lowercase + punct strip\nNLTK stopwords removed\nLemmatisation\n(WordNet)"),
    ("3. Feature Eng.",    "TF-IDF + LSA  (50d)\nLDA topics  (3d)\nSentence-BERT  (384d)\nall-MiniLM-L6-v2"),
]
row2 = [
    ("4. Clustering",      "K-Means++ on 3 reps\nHierarchical Ward\non SBERT\n+ optimal-k search"),
    ("5. Summarization",   "Extractive sentence\nranking via TF-IDF\nROUGE-1/2/L\nKeyword coverage"),
    ("6. Visualization",   "t-SNE / PCA 2D\nMetrics bar chart\nTopic heatmap\nClustering report"),
    ("7. Report + Demo",   "Self-contained HTML\n+ Streamlit web app\nGroup22 PPTX deck\nGitHub repo"),
]

# Geometry
TOP_Y     = 1.65
ROW_H     = 2.40
ROW_GAP   = 0.45            # vertical gap between row 1 and row 2
HEADER_H  = 0.50            # gold header strip height

def stage_box(x, y, w, h, num_title, body, fill=LIGHT_NAVY):
    rect(sl, x, y, w, h, fill=fill)
    rect(sl, x, y, w, HEADER_H, fill=GOLD)
    txt(sl, num_title, x + 0.08, y + 0.07, w - 0.16, HEADER_H - 0.1,
        size=13, bold=True, color=DEEP_NAVY, align=PP_ALIGN.CENTER)
    txt(sl, body, x + 0.15, y + HEADER_H + 0.12, w - 0.3, h - HEADER_H - 0.2,
        size=11, color=WHITE, align=PP_ALIGN.CENTER)

def arrow_right(x, y, w=0.45, color=NAVY):
    """Simple right-pointing arrow shape."""
    arrow = sl.shapes.add_shape(13, Inches(x), Inches(y), Inches(w), Inches(0.35))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()

def arrow_down(x, y, h=0.35, color=NAVY):
    arrow = sl.shapes.add_shape(14, Inches(x), Inches(y), Inches(0.35), Inches(h))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()

# Row 1: 3 boxes evenly spaced
n1 = len(row1)
gap1 = 0.55
total1 = (13.33 - 2 * 0.6 - (n1 - 1) * gap1)   # 0.6 left/right margin
box1_w = total1 / n1
for i, (title, body) in enumerate(row1):
    x = 0.6 + i * (box1_w + gap1)
    stage_box(x, TOP_Y, box1_w, ROW_H, title, body, fill=LIGHT_NAVY)
    if i < n1 - 1:
        ax = x + box1_w + (gap1 - 0.45) / 2
        ay = TOP_Y + ROW_H / 2 - 0.18
        arrow_right(ax, ay)

# Down arrow at end of row 1 → start of row 2 right side (we read row 2 left to right too)
y_after_row1 = TOP_Y + ROW_H
y_before_row2 = y_after_row1 + ROW_GAP
# Place a single down arrow in the middle of the slide between rows
arrow_down(13.33 / 2 - 0.18, y_after_row1 + 0.05, h=ROW_GAP - 0.1)

# Row 2: 4 boxes evenly spaced
n2 = len(row2)
gap2 = 0.4
total2 = (13.33 - 2 * 0.6 - (n2 - 1) * gap2)
box2_w = total2 / n2
for i, (title, body) in enumerate(row2):
    x = 0.6 + i * (box2_w + gap2)
    stage_box(x, y_before_row2, box2_w, ROW_H, title, body,
              fill=DEEP_NAVY)
    if i < n2 - 1:
        ax = x + box2_w + (gap2 - 0.45) / 2
        ay = y_before_row2 + ROW_H / 2 - 0.18
        arrow_right(ax, ay)

# Bottom caption
cap_y = y_before_row2 + ROW_H + 0.15
txt(sl, "Stages 1-3 build three parallel document representations.  "
        "Stages 4-7 cluster, summarise, visualise, and deliver.",
    0.6, cap_y, 12.1, 0.4, size=12, italic=True,
    color=RGBColor(0x55, 0x55, 0x55), align=PP_ALIGN.CENTER)


# ── Slide 4: Dataset ─────────────────────────────────────────────────────────

sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xf4, 0xf6, 0xf9))
header_bar(sl, "Dataset", "Real USPTO Patents via HuggingFace big_patent")
footer(sl)

txt(sl, "5,000 real USPTO patent abstracts streamed from the big_patent dataset (Sharma et al., 2019)",
    0.6, 1.6, 12.1, 0.45, size=15, color=BLUE, bold=True)

# Table
headers = ["CPC Section", "CPC Code", "Technology Domain", "Count"]
rows = [
    ["g", "G06N", "Machine Learning / Artificial Intelligence", "1,000"],
    ["h", "H01L", "Semiconductor / Electronics", "1,000"],
    ["a", "A61B", "Biotechnology / Medical Devices", "1,000"],
    ["b", "B60W", "Autonomous Vehicles / Transportation", "1,000"],
    ["f", "H02S", "Renewable Energy / Power Systems", "1,000"],
]
col_ws = [1.2, 1.3, 6.5, 1.2]
col_xs = [0.6, 1.85, 3.2, 9.75]
row_h  = 0.45
start_y = 2.2

# header row
for j, (hdr, cx, cw) in enumerate(zip(headers, col_xs, col_ws)):
    rect(sl, cx, start_y, cw, row_h, fill=NAVY)
    txt(sl, hdr, cx + 0.05, start_y + 0.07, cw - 0.1, row_h - 0.1,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i, row in enumerate(rows):
    y = start_y + (i + 1) * row_h
    bg = LIGHT if i % 2 == 0 else WHITE
    for j, (val, cx, cw) in enumerate(zip(row, col_xs, col_ws)):
        rect(sl, cx, y, cw, row_h, fill=bg, line=RGBColor(0xcc, 0xcc, 0xcc))
        align = PP_ALIGN.CENTER if j != 2 else PP_ALIGN.LEFT
        txt(sl, val, cx + 0.06, y + 0.07, cw - 0.12, row_h - 0.1,
            size=12, color=BLUE, align=align)

txt(sl, "Total: 5,000 patents  |  5 technology domains  |  1,000 patents per domain",
    0.6, 4.65, 12.1, 0.4, size=14, bold=True, color=NAVY)

bullet_block(sl, [
    "Abstracts range 40-2000 words; titles derived from first sentence of abstract",
    "Streamed directly from HuggingFace — no local download required",
    "Saved to DATA/real_patents.csv in the repository",
], 0.6, 5.15, 12.1, 1.6, size=13, color=RGBColor(0x33, 0x33, 0x33))


# ── Slide 5: Feature Engineering ─────────────────────────────────────────────

sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xf4, 0xf6, 0xf9))
header_bar(sl, "Feature Engineering", "Three parallel document representations")
footer(sl)

panels = [
    ("TF-IDF + LSA", NAVY,
     ["Vocabulary: 6,000 top unigrams",
      "TF-IDF matrix: 5,000 x 6,000",
      "SVD reduction -> 50 latent dimensions",
      "Captures lexical frequency patterns",
      "Fast, interpretable, no GPU needed"]),
    ("LDA Topic Model", BLUE,
     ["Optimal topics: 3 (coherence = 0.071)",
      "Document-topic distribution: 5,000 x 3",
      "Dirichlet prior: alpha=0.1, beta=0.01",
      "Discovers latent semantic themes",
      "Best clustering silhouette: 0.762"]),
    ("Sentence-BERT", RGBColor(0x10, 0x50, 0x80),
     ["Model: all-MiniLM-L6-v2 (HuggingFace)",
      "Dense embeddings: 5,000 x 384 dims",
      "Semantic similarity in cosine space",
      "Context-aware, handles paraphrases",
      "Embedding cache: skip re-encoding"]),
]
panel_w = 3.9
for i, (title, color, bullets) in enumerate(panels):
    x = 0.3 + i * (panel_w + 0.27)
    rect(sl, x, 1.55, panel_w, 5.3, fill=color)
    txt(sl, title, x + 0.1, 1.6, panel_w - 0.2, 0.55,
        size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    rect(sl, x + 0.1, 2.2, panel_w - 0.2, 0.03, fill=GOLD)
    bullet_block(sl, bullets, x + 0.2, 2.35, panel_w - 0.4, 4.2,
                 size=13, color=WHITE, bullet="-")


# ── Slide 6: Clustering Results ───────────────────────────────────────────────

sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xf4, 0xf6, 0xf9))
header_bar(sl, "Clustering Evaluation", "Silhouette up, Davies-Bouldin down, Stability up = better")
footer(sl)

methods = [
    ("TF-IDF + KMeans",      0.113, 3.603, 0.908),
    ("LDA + KMeans",         0.762, 0.618, 0.999),
    ("SBERT + KMeans",       0.045, 5.061, 0.940),
    ("SBERT + Hierarchical", 0.026, 6.010, 0.940),
]

hdrs = ["Method", "Silhouette (up)", "Davies-Bouldin (down)", "Stability ARI (up)", "Best?"]
col_ws2 = [3.5, 2.2, 2.5, 2.4, 1.3]
col_xs2 = [0.4, 3.95, 6.2, 8.75, 11.2]
row_h2  = 0.52
sy = 1.6

for j, (hdr, cx, cw) in enumerate(zip(hdrs, col_xs2, col_ws2)):
    rect(sl, cx, sy, cw, row_h2, fill=NAVY)
    txt(sl, hdr, cx + 0.05, sy + 0.1, cw - 0.1, row_h2 - 0.1,
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i, (method, sil, db, stab) in enumerate(methods):
    y = sy + (i + 1) * row_h2
    best = (sil == max(r[1] for r in methods))
    bg = RGBColor(0xd4, 0xed, 0xda) if best else (LIGHT if i % 2 == 0 else WHITE)
    vals = [method, f"{sil:.3f}", f"{db:.3f}", f"{stab:.3f}", "BEST" if best else "--"]
    for j, (val, cx, cw) in enumerate(zip(vals, col_xs2, col_ws2)):
        rect(sl, cx, y, cw, row_h2, fill=bg, line=RGBColor(0xcc, 0xcc, 0xcc))
        txt(sl, val, cx + 0.06, y + 0.1, cw - 0.12, row_h2 - 0.1,
            size=13, bold=(j == 4 and best), color=NAVY,
            align=(PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER))

txt(sl, "Key Findings:", 0.4, 4.7, 12.5, 0.35, size=14, bold=True, color=NAVY)
bullet_block(sl, [
    "LDA + KMeans achieves the highest Silhouette (0.762) and near-perfect Stability (0.999) on 5,000 real patents",
    "TF-IDF + KMeans yields the strongest interpretable clusters (low DB = 3.60, k=5 matches CPC sections)",
    "SBERT embeddings show lower silhouette typical of real-world dense semantic spaces with overlapping domains",
], 0.4, 5.1, 12.5, 1.7, size=13, color=RGBColor(0x22, 0x22, 0x22))


# ── Slide 7: Visualizations ───────────────────────────────────────────────────

sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xf4, 0xf6, 0xf9))
header_bar(sl, "Cluster Visualizations", "t-SNE 2D projections of all four clustering methods")
footer(sl)

# all_projections_tsne is 3.92:1 — fit it preserving that wide ratio
fit_picture(sl, os.path.join(OUT_DIR, "all_projections_tsne.png"),
            l=0.4, t=2.4, max_w=12.53, max_h=4.0)
txt(sl, "Each panel shows a 2D t-SNE projection of patents, colored by cluster.  "
        "TF-IDF and LDA produce visibly tighter, more separable clusters than SBERT on this corpus.",
    0.4, 6.55, 12.53, 0.6, size=12, italic=True,
    color=RGBColor(0x55, 0x55, 0x55), align=PP_ALIGN.CENTER)


# ── Slide 8: Metrics Comparison Chart ────────────────────────────────────────

sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xf4, 0xf6, 0xf9))
header_bar(sl, "Metrics Comparison", "Side-by-side bar chart: Silhouette, Davies-Bouldin, Stability ARI")
footer(sl)

# metrics_comparison is 2.02:1 — fit cleanly without stretch
fit_picture(sl, os.path.join(OUT_DIR, "metrics_comparison.png"),
            l=0.7, t=1.6, max_w=11.93, max_h=5.3)


# ── Slide 9: Summarization Evaluation ────────────────────────────────────────

sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xf4, 0xf6, 0xf9))
header_bar(sl, "Summarization Evaluation", "Extractive TF-IDF sentence ranking + ROUGE metrics")
footer(sl)

sum_metrics = [
    ("TF-IDF + KMeans",      0.260, 0.006, 0.193, 0.125, 0.025),
    ("LDA + KMeans",         0.280, 0.040, 0.239, 0.200, 0.039),
    ("SBERT + KMeans",       0.266, 0.013, 0.206, 0.100, 0.028),
    ("SBERT + Hierarchical", 0.284, 0.034, 0.203, 0.100, 0.033),
]
hdrs3 = ["Method", "ROUGE-1", "ROUGE-2", "ROUGE-L", "Kw Coverage", "Centroid Sim"]
col_ws3 = [3.3, 1.8, 1.8, 1.8, 2.0, 2.0]
col_xs3 = [0.4, 3.75, 5.6, 7.45, 9.3, 11.35]
row_h3  = 0.52
sy3 = 1.6

for j, (hdr, cx, cw) in enumerate(zip(hdrs3, col_xs3, col_ws3)):
    rect(sl, cx, sy3, cw, row_h3, fill=NAVY)
    txt(sl, hdr, cx + 0.05, sy3 + 0.1, cw - 0.1, row_h3 - 0.1,
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i, (method, r1, r2, rl, cov, prox) in enumerate(sum_metrics):
    y = sy3 + (i + 1) * row_h3
    bg = LIGHT if i % 2 == 0 else WHITE
    vals = [method, f"{r1:.3f}", f"{r2:.3f}", f"{rl:.3f}", f"{cov:.3f}", f"{prox:.3f}"]
    for j, (val, cx, cw) in enumerate(zip(vals, col_xs3, col_ws3)):
        rect(sl, cx, y, cw, row_h3, fill=bg, line=RGBColor(0xcc, 0xcc, 0xcc))
        txt(sl, val, cx + 0.06, y + 0.1, cw - 0.12, row_h3 - 0.1,
            size=13, color=NAVY,
            align=(PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER))

txt(sl, "How summarization works:", 0.4, 4.55, 12.5, 0.35, size=14, bold=True, color=NAVY)
bullet_block(sl, [
    "Each cluster's documents are merged; sentences ranked by TF-IDF cosine similarity to centroid",
    "Top-3 sentences selected as the cluster summary (order preserved for readability)",
    "ROUGE-1 measures unigram overlap; ROUGE-L measures longest common subsequence with representative titles",
    "LDA + KMeans wins ROUGE-2 (0.040), ROUGE-L (0.239) and Coverage (0.200) -- best balance overall",
], 0.4, 4.95, 12.5, 1.9, size=13, color=RGBColor(0x22, 0x22, 0x22))


# ── Slide 10: LDA Topic Distribution ─────────────────────────────────────────

sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xf4, 0xf6, 0xf9))
header_bar(sl, "LDA Topic Modeling", "Optimal coherence at 3 topics (coherence = 0.071)")
footer(sl)

# lda_topics ratio 1.33:1 — keep it on the left half
# doc_topic_heatmap ratio 1.06:1 (near square) — right half
fit_picture(sl, os.path.join(OUT_DIR, "lda_topics.png"),
            l=0.3, t=1.6, max_w=6.4, max_h=5.2)
fit_picture(sl, os.path.join(OUT_DIR, "doc_topic_heatmap.png"),
            l=6.9, t=1.6, max_w=6.1, max_h=5.2)

txt(sl, "Top words per topic (left)  ·  Document-topic probability heatmap (right)",
    0.4, 6.95, 12.5, 0.4, size=12, italic=True,
    color=RGBColor(0x55, 0x55, 0x55), align=PP_ALIGN.CENTER)


# ── Slide 11: Optimal k Selection ────────────────────────────────────────────

sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xf4, 0xf6, 0xf9))
header_bar(sl, "Optimal k Selection", "Silhouette & Davies-Bouldin vs k (SBERT embeddings)")
footer(sl)

# optimal_k_SBERT ratio 2.04:1 — fit to ~10 inch wide
fit_picture(sl, os.path.join(OUT_DIR, "optimal_k_SBERT.png"),
            l=1.5, t=1.6, max_w=10.33, max_h=4.6)

bullet_block(sl, [
    "k evaluated from 2 to 9 on SBERT embeddings using Silhouette and Davies-Bouldin",
    "Diagnostic optimal k = 2 by silhouette, but k = 5 chosen for clustering to match the 5 CPC sections",
    "Low SBERT silhouette values are expected: real patents have nuanced, overlapping technical language",
], 0.5, 6.4, 12.3, 0.85, size=12, color=RGBColor(0x22, 0x22, 0x22))


# ── Slide 12: Key Results Summary ────────────────────────────────────────────

sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xf4, 0xf6, 0xf9))
header_bar(sl, "Key Results & Contributions", "What we built and what we found")
footer(sl)

panels2 = [
    ("Pipeline Contributions", NAVY, [
        "End-to-end pipeline: ingestion -> preprocessing -> features -> clustering -> summarization -> report",
        "Three feature representations compared on 5,000 real USPTO patents",
        "Automated optimal-k and optimal-topics search with coherence scoring",
        "SBERT embedding cache + memory-aware feature builder for scale",
        "Self-contained HTML report + interactive Streamlit demo app",
    ]),
    ("Evaluation Findings", BLUE, [
        "LDA + KMeans: best overall (Sil=0.762, Stab=0.999, ROUGE-L=0.239)",
        "TF-IDF + KMeans: cleanest interpretable clusters (DB=3.60)",
        "SBERT + Hierarchical: best ROUGE-1 (0.284) -- captures paraphrase",
        "Pipeline scaled 10x from 500 to 5,000 patents in ~6 minutes",
        "Stability ARI > 0.90 across all four methods at this scale",
    ]),
]
pw = 6.0
for i, (title, color, bullets) in enumerate(panels2):
    x = 0.3 + i * (pw + 0.73)
    rect(sl, x, 1.55, pw, 5.3, fill=color)
    txt(sl, title, x + 0.1, 1.6, pw - 0.2, 0.55,
        size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    rect(sl, x + 0.1, 2.2, pw - 0.2, 0.03, fill=GOLD)
    bullet_block(sl, bullets, x + 0.2, 2.35, pw - 0.35, 4.3,
                 size=13, color=WHITE, bullet="-")


# ── Slide 13: GitHub & Thank You ─────────────────────────────────────────────

sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
rect(sl, 0, 0, 13.33, 0.08, fill=GOLD)
rect(sl, 0, 7.42, 13.33, 0.08, fill=GOLD)

txt(sl, "Thank You", 1, 1.0, 11.33, 1.0, size=40, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)
txt(sl, "Questions?", 1, 2.0, 11.33, 0.6, size=22, color=GOLD, align=PP_ALIGN.CENTER)

rect(sl, 2.5, 2.9, 8.33, 0.04, fill=GOLD)

txt(sl, "GitHub Repository", 1, 3.1, 11.33, 0.5,
    size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txt(sl, "github.com/akshith-22/PatentSummarizer",
    1, 3.55, 11.33, 0.5, size=15, color=RGBColor(0xa8, 0xd8, 0xea), align=PP_ALIGN.CENTER)

rect(sl, 2.5, 4.2, 8.33, 0.04, fill=GOLD)

info = [
    "Group 22  |  Project 9  |  CSE 573 Spring 2026",
    "Arizona State University",
    "Akshith Reddy Vempati  |  Reethika Gogireddy",
]
for i, line in enumerate(info):
    txt(sl, line, 1, 4.35 + i * 0.55, 11.33, 0.5,
        size=15, color=WHITE, align=PP_ALIGN.CENTER)

txt(sl, "Pipeline: TF-IDF + LDA + Sentence-BERT  ->  K-Means + Hierarchical Clustering  ->  ROUGE Evaluation",
    1, 6.4, 11.33, 0.5, size=12, italic=True,
    color=RGBColor(0xa8, 0xd8, 0xea), align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────

out = os.path.join(os.path.dirname(__file__),
                   "Group22-Project9-SP26-Group-DEMO-Presentation.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
